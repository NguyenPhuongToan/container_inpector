import re
from collections.abc import Sequence
from datetime import datetime
from pathlib import Path
from typing import Any

from app.core.config import settings
from app.core.constants import REPORT_ROOT
from app.database.models import Inspection


class FittingPhotoExportError(RuntimeError):
    """Raised when the fitting photo PowerPoint cannot be generated."""


OLD_CONTAINER_NUMBERS = [
    "WHSU 0008780",
    "WHSU 0514557",
    "WHSU 0171761",
    "WHSU 0305351",
    "FYCU 7158303",
    "FYC U 7158303",
]
OLD_FLEXITANK_NUMBERS = [
    "23TT3-2601-005",
    "23TT3-2601-012",
    "23TT3-2601-004",
    "23TT3-2601-006",
    "23TT3-2601-009",
]


def _import_pptx() -> Any:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise FittingPhotoExportError(
            "python-pptx is not installed. Install backend requirements before exporting PPTX."
        ) from exc

    return Presentation, MSO_SHAPE_TYPE


def _report_dir(booking_number: str) -> Path:
    safe_booking = re.sub(r"[^A-Za-z0-9_.-]+", "_", booking_number).strip("_") or "booking"
    report_dir = REPORT_ROOT / "fitting_photos" / safe_booking
    report_dir.mkdir(parents=True, exist_ok=True)
    return report_dir


def _template_path() -> Path:
    path = Path(settings.fitting_photo_template_path)
    if not path.is_absolute():
        path = Path.cwd() / path

    if not path.exists():
        raise FittingPhotoExportError(
            f"Fitting photo template was not found at {path}."
        )

    return path


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return shape.text or ""


def _replace_in_text_frame(shape: Any, replacements: dict[str, str]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return

    normalized_shape_text = re.sub(r"\s+", " ", shape.text or "").strip()
    for old, new in replacements.items():
        if normalized_shape_text == old:
            shape.text = new
            return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            text = run.text
            for old, new in replacements.items():
                text = text.replace(old, new)
            run.text = text


def _remove_slides(presentation: Any, indexes: Sequence[int]) -> None:
    slide_id_list = presentation.slides._sldIdLst  # noqa: SLF001
    for index in sorted(indexes, reverse=True):
        slide_id = slide_id_list[index]
        relationship_id = slide_id.rId
        presentation.part.drop_rel(relationship_id)
        slide_id_list.remove(slide_id)


def _replace_slide_pictures(slide: Any, image_paths: list[Path], picture_shape_type: Any) -> None:
    pictures = [
        shape
        for shape in slide.shapes
        if shape.shape_type == picture_shape_type.PICTURE
    ]
    pictures.sort(key=lambda shape: (shape.top, shape.left))

    for shape, image_path in zip(pictures, image_paths, strict=False):
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        element = shape._element  # noqa: SLF001
        element.getparent().remove(element)
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)


def _image_path(image: Any) -> Path | None:
    path = Path(image.path)
    if not path.exists():
        return None
    if path.suffix.lower() not in {".jpg", ".jpeg", ".png", ".webp"}:
        return None
    return path


def _inspection_image_paths(inspection: Inspection) -> list[Path]:
    paths: list[Path] = []
    for image in sorted(inspection.images, key=lambda item: item.angle):
        image_path = _image_path(image)
        if image_path is not None:
            paths.append(image_path)
    return paths[:12]


def _friendly_date(value: datetime | None) -> str:
    value = value or datetime.now()
    return value.strftime("%d %B, %Y")


def _cover_replacements(booking_number: str, quantity: int, created_at: datetime | None) -> dict[str, str]:
    return {
        "039GX48205": booking_number,
        "Quantity: 5 pcs": f"Quantity: {quantity} pcs",
        "25 May,": _friendly_date(created_at).rsplit(" ", 1)[0] + ",",
        "2026": _friendly_date(created_at).rsplit(" ", 1)[1],
    }


def _fill_header(slide: Any, index: int, inspection: Inspection) -> None:
    replacements = {
        old_value: inspection.container_number
        for old_value in OLD_CONTAINER_NUMBERS
    }
    replacements.update(
        {
            old_value: inspection.flexitank_number or "-"
            for old_value in OLD_FLEXITANK_NUMBERS
        }
    )

    for shape in slide.shapes:
        _replace_in_text_frame(shape, replacements)

    if index == 1:
        return

    numbered_shapes = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and re.fullmatch(r"\s*\d+\s*\.?\s*", _shape_text(shape))
    ]
    if numbered_shapes:
        numbered_shapes.sort(key=lambda shape: (shape.top, shape.left))
        numbered_shapes[0].text = f"{index}."


def generate_fitting_photo_pptx(inspections: Sequence[Inspection]) -> Path:
    if not inspections:
        raise FittingPhotoExportError("No inspections were selected for the PPTX export.")

    if len(inspections) > 5:
        raise FittingPhotoExportError(
            "The current template supports up to 5 containers per booking export."
        )

    booking_number = inspections[0].booking_number
    normalized_booking = booking_number.strip().lower()
    for inspection in inspections:
        if inspection.booking_number.strip().lower() != normalized_booking:
            raise FittingPhotoExportError("All inspections in a fitting photo export must share one booking number.")

    Presentation, MSO_SHAPE_TYPE = _import_pptx()
    presentation = Presentation(str(_template_path()))

    block_count = len(inspections)
    required_slide_count = 1 + (block_count * 3)
    if len(presentation.slides) < required_slide_count:
        raise FittingPhotoExportError("The fitting photo template does not contain enough container slide blocks.")

    _remove_slides(
        presentation,
        range(required_slide_count, len(presentation.slides)),
    )

    cover_created_at = min((inspection.created_at for inspection in inspections), default=None)
    for shape in presentation.slides[0].shapes:
        _replace_in_text_frame(
            shape,
            _cover_replacements(booking_number, len(inspections), cover_created_at),
        )

    for inspection_index, inspection in enumerate(inspections[:block_count], start=1):
        image_paths = _inspection_image_paths(inspection)
        if len(image_paths) < 12:
            raise FittingPhotoExportError(
                f"Inspection {inspection.container_number} needs 12 saved photos for PPTX export."
            )

        block_start = 1 + ((inspection_index - 1) * 3)
        _fill_header(presentation.slides[block_start], inspection_index, inspection)
        for offset in range(3):
            slide = presentation.slides[block_start + offset]
            _replace_slide_pictures(
                slide,
                image_paths[offset * 4:(offset + 1) * 4],
                MSO_SHAPE_TYPE,
            )

    report_dir = _report_dir(booking_number)
    output_path = report_dir / f"fitting_photo_{re.sub(r'[^A-Za-z0-9_.-]+', '_', booking_number)}.pptx"
    presentation.save(output_path)
    return output_path
